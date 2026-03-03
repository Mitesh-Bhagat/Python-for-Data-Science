from pptx import Presentation
from pptx.util import Inches, Pt

def add_slide(prs, title_text, content_text):
    # Use layout 1 (Title and Content)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set Title
    title = slide.shapes.title
    title.text = title_text
    
    # Set Content
    content = slide.placeholders[1]
    content.text = content_text

# Initialize Presentation
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0] # Title Slide Layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "TechCorp Global IAM Implementation Strategy"
subtitle.text = "Securing the Future: A Roadmap for Seamless Digital Transformation\nDecember 2025"

# Slide 2: Strategic Alignment
add_slide(prs, "Strategic Alignment & Business Goals",
    "• Fortify Cybersecurity: Shift to Zero Trust architecture.\n"
    "• Streamline Operations: Automate user lifecycle to reduce helpdesk tickets by 40%.\n"
    "• Enable Agility: Seamless SSO access for cloud and legacy tools.\n"
    "• Ensure Compliance: Centralize audit trails for GDPR/SOX."
)

# Slide 3: Implementation Roadmap
add_slide(prs, "Implementation Roadmap (Phased)",
    "• Phase 1 (Foundation): Establish infrastructure, clean AD data, define RBAC.\n"
    "• Phase 2 (Integration): Connect HR Systems (Workday/SAP) & Critical SaaS (M365).\n"
    "• Phase 3 (Security): Deploy Identity Gateways (Legacy) & Risk-based MFA.\n"
    "• Phase 4 (Rollout): Global regional go-live & legacy decommissioning."
)

# Slide 4: Integration Challenges
add_slide(prs, "Overcoming Integration Challenges",
    "• Legacy Systems: Use Identity Proxies to inject headers (no code changes required).\n"
    "• Cloud Sprawl: Standardize on SAML 2.0/OIDC protocols for all new apps.\n"
    "• Third-Party Access: Vendor Access Portal with Just-in-Time (JIT) provisioning."
)

# Slide 5: Automated User Lifecycle (JML)",
add_slide(prs, "Automated User Lifecycle (JML)",
    "• Joiner: HR triggers auto-provisioning of birthright access on Day 1.\n"
    "• Mover: Role changes automatically adjust permissions (prevent creep).\n"
    "• Leaver: Termination signals immediate kill-switch across all systems."
)

# Slide 6: Resource & Governance
add_slide(prs, "Resource Requirements",
    "• Core Team: IAM Architect, Integration Specialists, Security Engineers.\n"
    "• Infrastructure: High-availability clusters across 3 global regions.\n"
    "• Governance: Bi-weekly Steering Committee for milestone sign-off."
)

# Slide 7: Risks
add_slide(prs, "Risk Management",
    "• Data Quality: Pre-implementation cleansing sprint for HR/AD data.\n"
    "• User Pushback: extensive training and 'invisible' adaptive MFA.\n"
    "• Downtime: Testing in staging and rollback snapshots for critical apps."
)

# Slide 8: Next Steps
add_slide(prs, "Conclusion & Next Steps",
    "1. Formal sign-off on Phase 1 Scope.\n"
    "2. Kick-off workshop with HR for data mapping.\n"
    "3. Provision staging environments.\n\n"
    "Ready to execute."
)

# Save the file
output_file = "TechCorp_IAM_Implementation.pptx"
prs.save(output_file)
print(f"Presentation saved successfully as {output_file}")